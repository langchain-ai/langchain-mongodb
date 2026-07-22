"""Shared pytest fixtures for unit, integration, and e2e tests."""

from __future__ import annotations

import io
import os
import textwrap
from pathlib import Path
from unittest.mock import MagicMock

import boto3
import pytest
from moto import mock_aws

# ---------------------------------------------------------------------------
# Fixture data paths
# ---------------------------------------------------------------------------

FIXTURES_DIR = Path(__file__).parent / "fixtures"


# ---------------------------------------------------------------------------
# mongomock <-> pymongo>=4.9 compatibility
#
# pymongo 4.9 added a `sort` kwarg to bulk update/replace operations and
# forwards it (as None) to the bulk builder. mongomock's BulkOperationBuilder
# does not accept it yet, so strip it when unset.
# ---------------------------------------------------------------------------


def _patch_mongomock_sort_kwarg() -> None:
    try:
        from mongomock.collection import BulkOperationBuilder
    except ImportError:
        return

    for method_name in ("add_update", "add_replace"):
        original = getattr(BulkOperationBuilder, method_name, None)
        if original is None or getattr(original, "_sort_kwarg_patched", False):
            continue

        def _make_wrapper(orig):
            def wrapper(self, *args, **kwargs):
                if kwargs.get("sort") is None:
                    kwargs.pop("sort", None)
                return orig(self, *args, **kwargs)

            wrapper._sort_kwarg_patched = True
            return wrapper

        setattr(BulkOperationBuilder, method_name, _make_wrapper(original))


_patch_mongomock_sort_kwarg()


# ---------------------------------------------------------------------------
# AWS / S3 fixtures (moto)
# ---------------------------------------------------------------------------


@pytest.fixture(scope="function")
def aws_credentials():
    """Fake AWS creds so moto doesn't hit real AWS."""
    os.environ["AWS_ACCESS_KEY_ID"] = "testing"
    os.environ["AWS_SECRET_ACCESS_KEY"] = "testing"
    os.environ["AWS_SECURITY_TOKEN"] = "testing"
    os.environ["AWS_SESSION_TOKEN"] = "testing"
    os.environ["AWS_DEFAULT_REGION"] = "us-east-1"
    yield
    for k in (
        "AWS_ACCESS_KEY_ID",
        "AWS_SECRET_ACCESS_KEY",
        "AWS_SECURITY_TOKEN",
        "AWS_SESSION_TOKEN",
        "AWS_DEFAULT_REGION",
    ):
        os.environ.pop(k, None)


@pytest.fixture(scope="function")
def s3_bucket(aws_credentials):
    """Create a moto-mocked S3 bucket and yield its name."""
    with mock_aws():
        client = boto3.client("s3", region_name="us-east-1")
        client.create_bucket(Bucket="test-bucket")
        yield "test-bucket"


@pytest.fixture(scope="function")
def s3_backend(s3_bucket):
    """Return an S3Backend pointed at the moto-mocked bucket."""
    from langchain_mongodb_deepagents_vfs.backends.s3 import S3Backend

    with mock_aws():
        yield S3Backend(bucket_name=s3_bucket, region_name="us-east-1")


@pytest.fixture(scope="function")
def sqs_queue(aws_credentials):
    """Create a moto-mocked SQS queue and yield its URL."""
    with mock_aws():
        client = boto3.client("sqs", region_name="us-east-1")
        response = client.create_queue(QueueName="test-queue")
        yield response["QueueUrl"]


# ---------------------------------------------------------------------------
# MongoDB fixtures (mongomock)
# ---------------------------------------------------------------------------


@pytest.fixture(scope="function")
def mongo_collection():
    """Return a mongomock Collection for the chunk store."""
    mongomock = pytest.importorskip("mongomock")
    client = mongomock.MongoClient()
    return client["langchain_mongodb_deepagents_vfs"]["chunks"]


# ---------------------------------------------------------------------------
# Embedder mock
# ---------------------------------------------------------------------------


@pytest.fixture
def mock_embedder():
    """Embedder that returns deterministic zero vectors."""
    from langchain_mongodb_deepagents_vfs.embedder import Embedder

    embedder = MagicMock(spec=Embedder)
    embedder._dimensions = 1024

    def _embed(chunks):
        return [[0.0] * 1024 for _ in chunks]

    embedder.embed_batch.side_effect = _embed
    return embedder


# ---------------------------------------------------------------------------
# Chunker fixture
# ---------------------------------------------------------------------------


@pytest.fixture
def chunker():
    from langchain_mongodb_deepagents_vfs.chunker import Chunker

    return Chunker(token_limit=128, overlap=16)


# ---------------------------------------------------------------------------
# Sample text fixture
# ---------------------------------------------------------------------------


@pytest.fixture
def sample_text_bytes() -> bytes:
    return textwrap.dedent("""\
        Line 1: The quick brown fox jumps over the lazy dog.
        Line 2: MongoDB Atlas provides vector search capabilities.
        Line 3: LangChain DeepAgents can use filesystem backends.
        Line 4: S3 is a reliable object store.
        Line 5: Chunking preserves context across splits.
    """).encode("utf-8")


@pytest.fixture(scope="session")
def sample_xlsx_bytes() -> bytes:
    """Two-sheet xlsx workbook with mixed string/number content."""
    openpyxl = pytest.importorskip("openpyxl")
    wb = openpyxl.Workbook()
    ws1 = wb.active
    ws1.title = "Revenue"
    ws1.append(["Region", "Q1", "Q2", "Q3", "Q4"])
    ws1.append(["EMEA", 100, 120, 130, 145])
    ws1.append(["APAC", 80, 95, 110, 130])
    ws1.append(["AMER", 200, 210, 225, 240])
    ws2 = wb.create_sheet("Notes")
    ws2.append(["Author", "Comment"])
    ws2.append(["alice", "Vector search outperforms keyword on paraphrased queries."])
    ws2.append(["bob", "Atlas $rankFusion combines BM25 and cosine via RRF."])
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


@pytest.fixture(scope="session")
def sample_xls_bytes() -> bytes:
    """Two-sheet legacy .xls workbook (BIFF8) via xlwt."""
    xlwt = pytest.importorskip("xlwt")
    wb = xlwt.Workbook()
    s1 = wb.add_sheet("Revenue")
    headers = ["Region", "Q1", "Q2", "Q3", "Q4"]
    rows = [
        ["EMEA", 100, 120, 130, 145],
        ["APAC", 80, 95, 110, 130],
        ["AMER", 200, 210, 225, 240],
    ]
    for col, h in enumerate(headers):
        s1.write(0, col, h)
    for r, row in enumerate(rows, start=1):
        for c, v in enumerate(row):
            s1.write(r, c, v)
    s2 = wb.add_sheet("Notes")
    s2.write(0, 0, "Author")
    s2.write(0, 1, "Comment")
    s2.write(1, 0, "alice")
    s2.write(1, 1, "Vector search outperforms keyword on paraphrased queries.")
    s2.write(2, 0, "bob")
    s2.write(2, 1, "Atlas rankFusion combines BM25 and cosine via RRF.")
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


@pytest.fixture(scope="session")
def sample_pptx_bytes() -> bytes:
    """Two-slide pptx with title + body text, a table, and speaker notes."""
    pptx_mod = pytest.importorskip("pptx")
    Presentation = pptx_mod.Presentation
    from pptx.util import Inches

    prs = Presentation()
    # Slide 1: title + content layout
    slide1 = prs.slides.add_slide(prs.slide_layouts[1])
    slide1.shapes.title.text = "Hybrid Retrieval Overview"
    body = slide1.placeholders[1].text_frame
    body.text = "BM25 captures lexical matches"
    body.add_paragraph().text = "Vector search captures semantic similarity"
    body.add_paragraph().text = "RRF fuses the two ranked lists"
    slide1.notes_slide.notes_text_frame.text = (
        "Mention $rankFusion server-side pipeline."
    )

    # Slide 2: title + a 2x2 table
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    slide2.shapes.title.text = "Benchmark Snapshot"
    tbl_shape = slide2.shapes.add_table(
        rows=2, cols=2, left=Inches(1), top=Inches(2), width=Inches(6), height=Inches(2)
    )
    table = tbl_shape.table
    table.cell(0, 0).text = "Method"
    table.cell(0, 1).text = "nDCG@10"
    table.cell(1, 0).text = "Hybrid (RRF)"
    table.cell(1, 1).text = "0.71"

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


@pytest.fixture
def sample_pdf_bytes() -> bytes:
    """Create a minimal valid PDF in-memory using pypdf if available."""
    try:
        from reportlab.pdfgen import canvas as rl_canvas

        buf = io.BytesIO()
        c = rl_canvas.Canvas(buf)
        c.drawString(100, 750, "Hello PDF World")
        c.save()
        return buf.getvalue()
    except ImportError:
        # Return a minimal hand-crafted PDF (1 page, empty text)
        return (
            b"%PDF-1.4\n1 0 obj<</Type /Catalog /Pages 2 0 R>>endobj\n"
            b"2 0 obj<</Type /Pages /Kids [3 0 R] /Count 1>>endobj\n"
            b"3 0 obj<</Type /Page /Parent 2 0 R /MediaBox [0 0 612 792]>>endobj\n"
            b"xref\n0 4\n0000000000 65535 f \n0000000009 00000 n \n"
            b"0000000058 00000 n \n0000000115 00000 n \n"
            b"trailer<</Size 4 /Root 1 0 R>>\nstartxref\n190\n%%EOF"
        )
