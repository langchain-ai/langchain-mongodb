"""Document chunker: bytes → list[Chunk].

Supports plain text (.txt/.md/.rst/.csv), PDF (.pdf), Word (.docx),
Excel (.xlsx, .xls), and PowerPoint (.pptx, .ppt). Format is detected by
file extension first, then magic bytes as a fallback.

Chunking strategy (uniform across all formats):
  - Each format extractor returns a list of (page_number, text) tuples,
    where one "page" corresponds to a natural logical unit:
      * PDF        → one page per PDF page
      * DOCX       → single logical page (whole document)
      * XLSX / XLS → one page per worksheet
      * PPTX / PPT → one page per slide
      * Text/MD/…  → single logical page
  - Each page's text is then split independently with a token-aware
    recursive splitter (tiktoken cl100k_base), 512 tokens / 64-token
    overlap — preserving the same per-page chunking behaviour for every
    supported format so embeddings stay coherent and comparable.
  - Each chunk carries source_path, chunk_index, page_number, char_start,
    char_end, line_start, content so SearchRouter can return DeepAgents-
    compatible GrepMatch.line directly from stored metadata.
"""

from __future__ import annotations

import io
import logging
import zipfile
from collections.abc import Sequence
from pathlib import PurePosixPath

import tiktoken

from langchain_mongodb_deepagents_vfs.dtypes import Chunk
from langchain_mongodb_deepagents_vfs.errors import AdapterError, ErrorCode

logger = logging.getLogger(__name__)

_ENCODING = "cl100k_base"
_DEFAULT_TOKEN_LIMIT = 512
_DEFAULT_OVERLAP = 64

_SUPPORTED_EXTENSIONS = {
    ".txt",
    ".md",
    ".rst",
    ".csv",
    ".pdf",
    ".docx",
    ".xlsx",
    ".xls",
    ".pptx",
    ".ppt",
}

# Magic byte signature for OLE2 Compound File Binary (legacy .xls / .ppt / .doc)
_OLE_MAGIC = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"

# Zip-bomb guardrails for OOXML (docx/xlsx/pptx) parsing. A crafted ~42 KB
# archive can expand to gigabytes; the Office parsers decompress every entry, so
# we inspect the ZIP central directory first and refuse anything past these caps.
_MAX_OOXML_COMPRESSED_BYTES = 100 * 1024 * 1024  # raw archive on disk
_MAX_OOXML_UNCOMPRESSED_BYTES = 500 * 1024 * 1024  # total inflated size
_MAX_OOXML_ENTRIES = 10_000  # entry count (guards zip-of-many-files)
_MAX_OOXML_EXPANSION_RATIO = 200  # uncompressed / compressed


class Chunker:
    """Converts raw bytes into a list of Chunk objects.

    Args:
        token_limit: Maximum tokens per chunk (default 512).
        overlap: Token overlap between consecutive chunks (default 64).
    """

    def __init__(
        self,
        token_limit: int = _DEFAULT_TOKEN_LIMIT,
        overlap: int = _DEFAULT_OVERLAP,
    ) -> None:
        self._token_limit = token_limit
        self._overlap = overlap
        self._enc = tiktoken.get_encoding(_ENCODING)

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    def chunk(self, path: str, data: bytes) -> list[Chunk]:
        """Extract text from *data* and split into overlapping chunks.

        Args:
            path: Source object key / path (used to derive filename + extension).
            data: Raw bytes of the object.

        Returns:
            Ordered list of Chunk objects.

        Raises:
            AdapterError(E9003): Unsupported file format.
            AdapterError(E9002): Extraction failure.
        """
        ext = PurePosixPath(path).suffix.lower()
        try:
            pages = self._extract(data, ext)
        except AdapterError:
            raise
        except Exception as exc:
            raise AdapterError(ErrorCode.E9002_CHUNKER_FAILED, str(exc)) from exc

        chunks: list[Chunk] = []
        chunk_index = 0
        for page_number, page_text in pages:
            for chunk_text, char_start, char_end, line_start in self._split(page_text):
                chunks.append(
                    Chunk(
                        source_path=path,
                        chunk_index=chunk_index,
                        content=chunk_text,
                        page_number=page_number,
                        char_start=char_start,
                        char_end=char_end,
                        line_start=line_start,
                    )
                )
                chunk_index += 1

        return chunks

    # ------------------------------------------------------------------
    # Format extractors
    # ------------------------------------------------------------------

    def _extract(self, data: bytes, ext: str) -> list[tuple[int, str]]:
        """Return list of (page_number, text) pairs."""
        if ext in (".txt", ".md", ".rst", ".csv"):
            return self._extract_text(data)
        if ext == ".pdf":
            return self._extract_pdf(data)
        if ext == ".docx":
            return self._extract_docx(data)
        if ext == ".xlsx":
            return self._extract_xlsx(data)
        if ext == ".xls":
            return self._extract_xls(data)
        if ext == ".pptx":
            return self._extract_pptx(data)
        if ext == ".ppt":
            return self._extract_ppt(data)

        # ---- magic-byte fallback ----
        if data[:4] == b"%PDF":
            return self._extract_pdf(data)
        if data[:2] == b"PK":
            # ZIP archive: could be docx / xlsx / pptx — peek inside
            kind = self._sniff_ooxml(data)
            if kind == "xlsx":
                return self._extract_xlsx(data)
            if kind == "pptx":
                return self._extract_pptx(data)
            if kind == "docx":
                return self._extract_docx(data)
            logger.warning(
                "Unrecognized ZIP archive '%s', attempting docx extraction.", ext
            )
            return self._extract_docx(data)
        if data[:8] == _OLE_MAGIC:
            # Legacy OLE2: try xls, then ppt
            try:
                return self._extract_xls(data)
            except AdapterError:
                return self._extract_ppt(data)
        # Try plain text as last resort; warn
        logger.warning("Unrecognized extension '%s', treating as plain text.", ext)
        return self._extract_text(data)

    @staticmethod
    def _guard_ooxml(data: bytes) -> None:
        """Reject OOXML zip bombs before handing bytes to an Office parser.

        Inspects the ZIP central directory (metadata only, no decompression) and
        raises AdapterError(E9002) if compressed size, total uncompressed size,
        entry count, or expansion ratio exceed safe bounds.
        """
        if len(data) > _MAX_OOXML_COMPRESSED_BYTES:
            raise AdapterError(
                ErrorCode.E9002_CHUNKER_FAILED,
                f"Archive is {len(data)} bytes, exceeds "
                f"{_MAX_OOXML_COMPRESSED_BYTES} byte limit",
            )
        try:
            with zipfile.ZipFile(io.BytesIO(data)) as zf:
                infos = zf.infolist()
        except zipfile.BadZipFile as exc:
            raise AdapterError(
                ErrorCode.E9002_CHUNKER_FAILED, f"Corrupt archive: {exc}"
            ) from exc

        if len(infos) > _MAX_OOXML_ENTRIES:
            raise AdapterError(
                ErrorCode.E9002_CHUNKER_FAILED,
                f"Archive has {len(infos)} entries, exceeds "
                f"{_MAX_OOXML_ENTRIES} limit",
            )
        total_uncompressed = sum(i.file_size for i in infos)
        total_compressed = sum(i.compress_size for i in infos)
        if total_uncompressed > _MAX_OOXML_UNCOMPRESSED_BYTES:
            raise AdapterError(
                ErrorCode.E9002_CHUNKER_FAILED,
                f"Archive inflates to {total_uncompressed} bytes, exceeds "
                f"{_MAX_OOXML_UNCOMPRESSED_BYTES} byte limit",
            )
        if (
            total_compressed > 0
            and total_uncompressed / total_compressed > _MAX_OOXML_EXPANSION_RATIO
        ):
            raise AdapterError(
                ErrorCode.E9002_CHUNKER_FAILED,
                f"Archive expansion ratio {total_uncompressed / total_compressed:.0f}x "
                f"exceeds {_MAX_OOXML_EXPANSION_RATIO}x limit",
            )

    @staticmethod
    def _sniff_ooxml(data: bytes) -> str | None:
        """Inspect a ZIP archive to identify which OOXML variant it is."""
        try:
            with zipfile.ZipFile(io.BytesIO(data)) as zf:
                names = zf.namelist()
        except zipfile.BadZipFile:
            return None
        if any(n.startswith("xl/") for n in names):
            return "xlsx"
        if any(n.startswith("ppt/") for n in names):
            return "pptx"
        if any(n.startswith("word/") for n in names):
            return "docx"
        return None

    def _extract_text(self, data: bytes) -> list[tuple[int, str]]:
        try:
            text = data.decode("utf-8", errors="replace")
        except Exception as exc:
            raise AdapterError(
                ErrorCode.E9002_CHUNKER_FAILED, f"UTF-8 decode failed: {exc}"
            ) from exc
        return [(0, text)]

    def _extract_pdf(self, data: bytes) -> list[tuple[int, str]]:
        try:
            from pypdf import PdfReader
        except ImportError as exc:
            raise AdapterError(
                ErrorCode.E9003_FORMAT_UNSUPPORTED, "pypdf is not installed"
            ) from exc

        # strict=False: tolerate malformed float tokens (common in reportlab/graphics-heavy PDFs)
        reader = PdfReader(io.BytesIO(data), strict=False)
        pages: list[tuple[int, str]] = []
        for i, page in enumerate(reader.pages):
            text = ""
            try:
                # layout mode preserves column order and table structure
                text = page.extract_text(extraction_mode="layout") or ""
            except Exception:
                pass
            if not text.strip():
                try:
                    # plain extraction as fallback for pages where layout mode yields nothing
                    text = page.extract_text() or ""
                except Exception:
                    text = ""
            pages.append((i, text))
        return pages

    def _extract_docx(self, data: bytes) -> list[tuple[int, str]]:
        try:
            import docx
        except ImportError as exc:
            raise AdapterError(
                ErrorCode.E9003_FORMAT_UNSUPPORTED, "python-docx is not installed"
            ) from exc
        self._guard_ooxml(data)
        doc = docx.Document(io.BytesIO(data))
        full_text = "\n".join(p.text for p in doc.paragraphs)
        return [(0, full_text)]

    def _extract_xlsx(self, data: bytes) -> list[tuple[int, str]]:
        """One (page_number, text) per worksheet. Rows rendered as TSV lines."""
        try:
            from openpyxl import load_workbook
        except ImportError as exc:
            raise AdapterError(
                ErrorCode.E9003_FORMAT_UNSUPPORTED,
                "openpyxl is not installed; pip install openpyxl",
            ) from exc
        self._guard_ooxml(data)
        wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        pages: list[tuple[int, str]] = []
        try:
            for page_number, sheet in enumerate(wb.worksheets):
                lines: list[str] = [f"# Sheet: {sheet.title}"]
                for row in sheet.iter_rows(values_only=True):
                    if row is None:
                        continue
                    cells = ["" if v is None else str(v) for v in row]
                    if any(c.strip() for c in cells):
                        lines.append("\t".join(cells))
                pages.append((page_number, "\n".join(lines)))
        finally:
            wb.close()
        return pages or [(0, "")]

    def _extract_xls(self, data: bytes) -> list[tuple[int, str]]:
        """Legacy .xls (BIFF) via xlrd. xlrd>=2.0 supports .xls only."""
        try:
            import xlrd
        except ImportError as exc:
            raise AdapterError(
                ErrorCode.E9003_FORMAT_UNSUPPORTED,
                "xlrd is not installed; pip install 'xlrd>=2.0.1'",
            ) from exc
        try:
            book = xlrd.open_workbook(file_contents=data)
        except Exception as exc:
            raise AdapterError(
                ErrorCode.E9002_CHUNKER_FAILED, f"xlrd failed: {exc}"
            ) from exc
        pages: list[tuple[int, str]] = []
        for page_number, sheet in enumerate(book.sheets()):
            lines: list[str] = [f"# Sheet: {sheet.name}"]
            for row_idx in range(sheet.nrows):
                cells = ["" if v is None else str(v) for v in sheet.row_values(row_idx)]
                if any(c.strip() for c in cells):
                    lines.append("\t".join(cells))
            pages.append((page_number, "\n".join(lines)))
        return pages or [(0, "")]

    def _extract_pptx(self, data: bytes) -> list[tuple[int, str]]:
        """One (page_number, text) per slide. Pulls text from shapes + notes."""
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise AdapterError(
                ErrorCode.E9003_FORMAT_UNSUPPORTED,
                "python-pptx is not installed; pip install python-pptx",
            ) from exc
        self._guard_ooxml(data)
        prs = Presentation(io.BytesIO(data))
        pages: list[tuple[int, str]] = []
        for page_number, slide in enumerate(prs.slides):
            parts: list[str] = []
            for shape in slide.shapes:
                # Text-bearing shapes
                if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = "".join(run.text for run in para.runs).strip()
                        if text:
                            parts.append(text)
                # Tables
                if getattr(shape, "has_table", False) and shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        if any(cells):
                            parts.append("\t".join(cells))
            # Speaker notes
            if slide.has_notes_slide:
                notes_tf = slide.notes_slide.notes_text_frame
                notes = (notes_tf.text or "").strip() if notes_tf is not None else ""
                if notes:
                    parts.append(f"[notes] {notes}")
            pages.append((page_number, "\n".join(parts)))
        return pages or [(0, "")]

    def _extract_ppt(self, data: bytes) -> list[tuple[int, str]]:
        """Legacy binary .ppt via olefile (best-effort text extraction).

        The binary PPT format has no clean pure-Python parser; we scan the
        ``PowerPoint Document`` OLE stream for embedded UTF-16LE / ASCII
        strings. Adequate for indexing, not for fidelity rendering.
        """
        try:
            import olefile
        except ImportError as exc:
            raise AdapterError(
                ErrorCode.E9003_FORMAT_UNSUPPORTED,
                "olefile is not installed; pip install olefile (required for legacy .ppt)",
            ) from exc
        try:
            ole = olefile.OleFileIO(io.BytesIO(data))
        except Exception as exc:
            raise AdapterError(
                ErrorCode.E9002_CHUNKER_FAILED, f"olefile failed: {exc}"
            ) from exc
        try:
            stream_name = "PowerPoint Document"
            if not ole.exists(stream_name):
                raise AdapterError(
                    ErrorCode.E9002_CHUNKER_FAILED,
                    f"Stream '{stream_name}' missing — not a valid .ppt file",
                )
            raw = ole.openstream(stream_name).read()
        finally:
            ole.close()

        text = self._scan_ole_strings(raw)
        return [(0, text)]

    @staticmethod
    def _scan_ole_strings(blob: bytes, min_len: int = 4) -> str:
        """Best-effort string extraction from a binary OLE stream."""
        out: list[str] = []
        # UTF-16LE runs: pairs of (printable ascii, 0x00)
        i = 0
        n = len(blob)
        while i < n - 1:
            run: list[int] = []
            j = i
            while j < n - 1 and 0x20 <= blob[j] <= 0x7E and blob[j + 1] == 0x00:
                run.append(blob[j])
                j += 2
            if len(run) >= min_len:
                out.append(bytes(run).decode("ascii", errors="replace"))
                i = j
            else:
                i += 1
        # ASCII runs
        ascii_run: list[int] = []
        for b in blob:
            if 0x20 <= b <= 0x7E or b in (0x09, 0x0A, 0x0D):
                ascii_run.append(b)
            else:
                if len(ascii_run) >= min_len:
                    out.append(bytes(ascii_run).decode("ascii", errors="replace"))
                ascii_run = []
        if len(ascii_run) >= min_len:
            out.append(bytes(ascii_run).decode("ascii", errors="replace"))
        return "\n".join(out)

    # ------------------------------------------------------------------
    # Token-aware recursive splitter
    # ------------------------------------------------------------------

    def _split(self, text: str) -> list[tuple[str, int, int, int]]:
        """Split *text* into overlapping token-bounded windows.

        Returns:
            List of (chunk_text, char_start, char_end, line_start).
        """
        tokens = self._enc.encode(text)
        results: list[tuple[str, int, int, int]] = []
        start = 0
        while start < len(tokens):
            end = min(start + self._token_limit, len(tokens))
            chunk_tokens = tokens[start:end]
            chunk_text = self._enc.decode(chunk_tokens)
            # Locate char positions in original text for metadata
            char_start = self._token_offset(text, tokens[:start])
            char_end = char_start + len(chunk_text)
            line_start = text[:char_start].count("\n")
            results.append((chunk_text, char_start, char_end, line_start))
            if end == len(tokens):
                break
            start += self._token_limit - self._overlap
        return results

    def _token_offset(self, text: str, prefix_tokens: Sequence[int]) -> int:
        """Return the character offset in *text* after decoding *prefix_tokens*."""
        if not prefix_tokens:
            return 0
        return len(self._enc.decode(list(prefix_tokens)))
